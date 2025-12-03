from typing import List
from pathlib import Path

from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader
from docx import Document as DocxDocument
from pptx import Presentation


SUPPORTED_EXTS = {".pdf", ".txt", ".md", ".docx", ".pptx"}


def load_documents(files: List[Path]) -> List[Document]:
    docs: List[Document] = []
    for file in files:
        ext = file.suffix.lower()
        if ext == ".pdf":
            docs.extend(_load_pdf(file))
        elif ext == ".docx":
            docs.extend(_load_docx(file))
        elif ext == ".pptx":
            docs.extend(_load_pptx(file))
        elif ext in (".txt", ".md"):
            docs.extend(_load_text(file))
        else:
            continue
    return docs


def _load_pdf(path: Path) -> List[Document]:
    loader = PyPDFLoader(str(path))
    raw_docs = loader.load()
    for d in raw_docs:
        d.metadata.update(
            {
                "source": str(path),
                "file_name": path.name,
                "type": "pdf",
                # PyPDFLoader already sets 'page' in metadata
            }
        )
    return raw_docs


def _load_docx(path: Path) -> List[Document]:
    doc = DocxDocument(str(path))
    chunks: List[Document] = []
    texts = []
    for para in doc.paragraphs:
        if para.text:
            texts.append(para.text)
    content = "\n".join(texts).strip()
    if content:
        chunks.append(
            Document(
                page_content=content,
                metadata={
                    "source": str(path),
                    "file_name": path.name,
                    "type": "docx",
                    "page": None,
                },
            )
        )
    return chunks


def _load_pptx(path: Path) -> List[Document]:
    prs = Presentation(str(path))
    docs: List[Document] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        content = "\n".join(texts).strip()
        if content:
            docs.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": str(path),
                        "file_name": path.name,
                        "type": "pptx",
                        "page": i,  # treat slide number as page
                    },
                )
            )
    return docs


def _load_text(path: Path) -> List[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    if not text.strip():
        return []
    return [
        Document(
            page_content=text,
            metadata={
                "source": str(path),
                "file_name": path.name,
                "type": path.suffix.lower().lstrip("."),
                "page": None,
            },
        )
    ]
